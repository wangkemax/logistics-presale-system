"""Solution presentation PPT generator.

Produces a 12-slide .pptx with:
- Cover, agenda, project overview, requirements summary
- Solution design, automation plan, cost analysis
- Financial indicators, risk summary, case references
- Implementation timeline, team, closing
"""

import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

import structlog

logger = structlog.get_logger()

# ── Style constants ──
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x05, 0x96, 0x69)


def _add_bg(slide, color=DARK_BLUE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_box(slide, text, left=Inches(0.8), top=Inches(0.5), width=Inches(8.4), font_size=28, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox


def _add_body_text(slide, text, left=Inches(0.8), top=Inches(1.5), width=Inches(8.4), height=Inches(4.5), font_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    return txBox


def _add_kpi_boxes(slide, kpis: list[tuple[str, str, str]], top=Inches(4)):
    """Add KPI indicator boxes in a row."""
    n = len(kpis)
    box_w = Inches(2)
    gap = Inches(0.3)
    total_w = n * box_w + (n - 1) * gap
    start_x = (Inches(10) - total_w) / 2

    for i, (label, value, subtitle) in enumerate(kpis):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(1, x, top, box_w, Inches(1.2))  # 1 = rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY_TEXT
        p2.alignment = PP_ALIGN.CENTER


def _add_table_slide(slide, title, headers, rows, top=Inches(1.5)):
    """Add a styled table to a slide."""
    _add_title_box(slide, title)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.8), top, Inches(8.4), Inches(0.4 * n_rows)
    )
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_TEXT


async def generate_solution_pptx(stage_outputs: dict, project_info: dict) -> bytes:
    """Generate a solution presentation PPT.

    Args:
        stage_outputs: All pipeline stage outputs.
        project_info: Project metadata.

    Returns:
        .pptx file as bytes.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    blank = prs.slide_layouts[6]  # blank layout

    solution = stage_outputs.get(5, {})
    automation = stage_outputs.get(6, {})
    cost_model = stage_outputs.get(8, {})
    risks = stage_outputs.get(9, {})
    benchmarks = stage_outputs.get(7, {})
    requirements = stage_outputs.get(1, {})

    # ═══ Slide 1: Cover ═══
    slide = prs.slides.add_slide(blank)
    _add_bg(slide, DARK_BLUE)

    _add_title_box(slide, project_info.get("name", "物流解决方案"),
                   top=Inches(1.5), font_size=32, color=WHITE)
    _add_body_text(slide, f"客户: {project_info.get('client_name', '')}\n"
                          f"行业: {project_info.get('industry', '')}\n"
                          f"日期: {datetime.now().strftime('%Y年%m月')}",
                   top=Inches(3), font_size=14)
    # Override text color to white
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = WHITE

    # ═══ Slide 2: Agenda ═══
    slide = prs.slides.add_slide(blank)
    _add_title_box(slide, "汇报议程")
    agenda = [
        "1. 项目概览与需求理解",
        "2. 物流解决方案设计",
        "3. 自动化技术方案",
        "4. 成本分析与报价",
        "5. 投资回报分析",
        "6. 风险管理",
        "7. 实施计划",
    ]
    _add_body_text(slide, "\n".join(agenda), font_size=16)

    # ═══ Slide 3: Project Overview ═══
    slide = prs.slides.add_slide(blank)
    _add_title_box(slide, "项目概览")
    overview = solution.get("executive_summary", "")
    key_metrics = requirements.get("key_metrics", {})
    metrics_text = ""
    if key_metrics:
        parts = []
        if key_metrics.get("warehouse_area_sqm"):
            parts.append(f"仓库面积: {key_metrics['warehouse_area_sqm']} ㎡")
        if key_metrics.get("daily_order_volume"):
            parts.append(f"日均订单: {key_metrics['daily_order_volume']} 单")
        if key_metrics.get("sku_count"):
            parts.append(f"SKU 数量: {key_metrics['sku_count']}")
        metrics_text = " | ".join(parts)
    _add_body_text(slide, f"{overview}\n\n{metrics_text}")

    # ═══ Slide 4: Requirements Summary ═══
    slide = prs.slides.add_slide(blank)
    reqs = requirements.get("requirements", [])
    p0_reqs = [r for r in reqs if r.get("priority") == "P0"][:8]
    if p0_reqs:
        _add_table_slide(slide, "核心需求 (P0)",
                         ["编号", "需求描述", "类别"],
                         [[r.get("id", ""), r.get("description", "")[:50], r.get("category", "")] for r in p0_reqs])
    else:
        _add_title_box(slide, "需求摘要")
        _add_body_text(slide, f"共提取 {len(reqs)} 项需求")

    # ═══ Slide 5: Solution Design ═══
    slide = prs.slides.add_slide(blank)
    _add_title_box(slide, "物流解决方案设计")
    wh = solution.get("warehouse_design", {})
    ops = solution.get("operations_design", {})
    sol_text = f"总面积: {wh.get('total_area_sqm', '——')} ㎡\n"
    sol_text += f"功能分区: {len(wh.get('zones', []))} 个区域\n"
    sol_text += f"流程设计: {wh.get('flow_design', '')}\n"
    picking = ops.get("picking", {})
    if isinstance(picking, dict):
        sol_text += f"\n拣选策略: {picking.get('strategy', '')}"
    _add_body_text(slide, sol_text)

    # ═══ Slide 6: Automation ═══
    slide = prs.slides.add_slide(blank)
    auto_recs = automation.get("recommendations", [])[:5]
    if auto_recs:
        _add_table_slide(slide, "自动化方案推荐",
                         ["技术", "应用场景", "评分", "投资(万)", "ROI"],
                         [
                             [r.get("technology", ""),
                              r.get("application_area", ""),
                              f"{r.get('suitability_score', 0)}/10",
                              f"{r.get('estimated_cost_cny', 0)/10000:.0f}",
                              f"{r.get('roi_percent', 0)}%"]
                             for r in auto_recs
                         ])
    else:
        _add_title_box(slide, "自动化方案")
        _add_body_text(slide, f"自动化等级: {automation.get('automation_level', '——')}")

    # ═══ Slide 7: Cost Breakdown ═══
    slide = prs.slides.add_slide(blank)
    breakdown = cost_model.get("cost_breakdown", {})
    cost_rows = []
    for key, label in [("labor", "人力"), ("facility", "场地"), ("equipment", "设备"), ("technology", "技术"), ("operations", "运营")]:
        cat = breakdown.get(key, {})
        y1 = cat.get("year1", 0) or 0
        y2 = cat.get("year2", 0) or 0
        y3 = cat.get("year3", 0) or 0
        cost_rows.append([label, f"¥{y1/10000:.0f}万", f"¥{y2/10000:.0f}万", f"¥{y3/10000:.0f}万"])

    if cost_rows:
        _add_table_slide(slide, "成本分析 (3 年)",
                         ["类别", "第1年", "第2年", "第3年"], cost_rows)
    else:
        _add_title_box(slide, "成本分析")
        _add_body_text(slide, "详见报价文件")

    # ═══ Slide 8: Financial Indicators ═══
    slide = prs.slides.add_slide(blank)
    _add_title_box(slide, "投资回报分析")
    indicators = cost_model.get("financial_indicators", {})
    _add_kpi_boxes(slide, [
        ("投资回报率", f"{indicators.get('roi_percent', 0):.1f}%", "ROI"),
        ("内部收益率", f"{indicators.get('irr_percent', 0):.1f}%", "IRR"),
        ("净现值", f"¥{indicators.get('npv_at_8pct', 0)/10000:.0f}万", "NPV @8%"),
        ("回本周期", f"{indicators.get('payback_months', 0)}个月", "Payback"),
    ], top=Inches(2))

    # ═══ Slide 9: Risks ═══
    slide = prs.slides.add_slide(blank)
    risk_items = risks.get("risk_matrix", [])[:6]
    if risk_items:
        _add_table_slide(slide, "风险管理",
                         ["风险", "可能性", "影响", "缓解措施"],
                         [[r.get("description", "")[:35], r.get("likelihood", ""),
                           r.get("impact", ""), r.get("mitigation", "")[:35]] for r in risk_items])
    else:
        _add_title_box(slide, "风险管理")
        _add_body_text(slide, f"整体风险等级: {risks.get('overall_risk_level', '——')}")

    # ═══ Slide 10: Case References ═══
    slide = prs.slides.add_slide(blank)
    cases = benchmarks.get("matched_cases", [])[:4]
    if cases:
        _add_table_slide(slide, "成功案例参考",
                         ["案例", "行业", "相似度", "适用性"],
                         [[c.get("case_name", ""), c.get("client_industry", ""),
                           f"{c.get('similarity_score', 0):.0%}",
                           c.get("applicable_to_current", "")[:30]] for c in cases])
    else:
        _add_title_box(slide, "案例参考")
        _add_body_text(slide, "详见标书附录")

    # ═══ Slide 11: Implementation Plan ═══
    slide = prs.slides.add_slide(blank)
    _add_title_box(slide, "实施计划")
    phases = [
        ["Phase 1", "项目启动与团队组建", "第 1-2 周"],
        ["Phase 2", "场地准备与设备采购", "第 3-8 周"],
        ["Phase 3", "系统部署与集成", "第 6-10 周"],
        ["Phase 4", "试运行与调优", "第 10-12 周"],
        ["Phase 5", "正式运营", "第 12 周起"],
    ]
    _add_table_slide(slide, "实施计划",
                     ["阶段", "主要工作", "时间"], phases, top=Inches(1.5))

    # ═══ Slide 12: Closing ═══
    slide = prs.slides.add_slide(blank)
    _add_bg(slide, DARK_BLUE)
    _add_title_box(slide, "感谢您的时间", top=Inches(1.8), font_size=32, color=WHITE)
    _add_body_text(slide, "期待与您的合作\n如有任何问题，请随时联系我们",
                   top=Inches(3), font_size=16)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER

    # Save
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
